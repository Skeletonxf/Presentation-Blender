from pptx import Presentation
import json
import copy
import six
import sys
import subprocess

json_data = json.loads(sys.argv[1])
# get input key from dictionary
data_list = json_data.get('input')

print("running")

# a slice of a presentation
class Slice:
    def __init__(self, data):
        self.filename = data.get('fp')
        self.presentation = Presentation(self.filename)
        length = len(self.presentation.slides)
        self.start = int(data.get('s'))
        self.end = int(data.get('e'))
        if self.end == 0:
            self.end = length
        if length < self.end:
            raise ValueError('Slice created with end slide number > than number of slides')

    def __str__(self):
        return self.filename + ": " + str(self.start) + " to " + str(self.end)


print("got to here")

slices = []
for data in data_list:
    slices.append(Slice(data))

def see(slices):
    for s in slices:
        print(s)

print("inputs")
see(slices)

def delete_slides(presentation, index):
        xml_slides = presentation.slides._sldIdLst
        slides = list(xml_slides)
        xml_slides.remove(slides[index])

command = 'libreoffice --show '
params = ''
for i in range(len(slices)):
    tmp = Presentation(slices[i].filename)
    c = 0
    for cntr in range(len(slices[i].presentation.slides)):
        if cntr < slices[i].start-c-1:
            delete_slides(tmp, cntr-c)
            c += 1
        elif cntr > slices[i].end-c-1:
            delete_slides(tmp, cntr-c)
            c += 1

    filename = 'presentation' + str(i) + '.pptx'
    print('Saving: ' + filename)
    tmp.save(filename)
    params = filename + ' ' + params

command += params
print('running ' + command)

process = subprocess.Popen(command.split(), stdout=subprocess.PIPE)
output, error = process.communicate()


